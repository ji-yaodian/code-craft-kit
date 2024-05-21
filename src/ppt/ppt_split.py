import os
from pptx import Presentation


def split_ppt(input_pptx_path):
    """
    将PPT文件拆分成单独的幻灯片文件。

    Args:
        input_pptx_path (str): 输入PPT文件路径。
    """

    # 创建输出目录
    output_pptx_dir = os.path.join(os.path.dirname(input_pptx_path), os.path.basename(input_pptx_path) + "_split_")
    if not os.path.exists(output_pptx_dir):
        os.makedirs(output_pptx_dir)

    # 读取PPT文件
    presentation = Presentation(input_pptx_path)

    # 获取幻灯片总数
    total_slides = len(presentation.slides)

    # 逐个拆分幻灯片
    for current_slide_index in range(total_slides):
        # 创建临时演示文稿
        temp_presentation = Presentation(input_pptx_path)

        # 删除当前幻灯片之前的幻灯片
        index = 0
        while index < current_slide_index:
            slide_id = temp_presentation.slides._sldIdLst[0].rId
            temp_presentation.part.drop_rel(slide_id)
            del temp_presentation.slides._sldIdLst[0]
            index += 1

        # 保留当前幻灯片
        index += 1

        # 删除当前幻灯片之后的幻灯片
        while index < total_slides:
            slide_id = temp_presentation.slides._sldIdLst[1].rId
            temp_presentation.part.drop_rel(slide_id)
            del temp_presentation.slides._sldIdLst[1]
            index += 1

        # 保存拆分后的幻灯片
        print(f"保存第 {current_slide_index + 1} 页幻灯片")
        temp_presentation.save(os.path.join(output_pptx_dir, f"{current_slide_index + 1}.pptx"))


if __name__ == "__main__":
    import sys

    # 输入PPT文件路径
    input_pptx_path = sys.argv[1]

    split_ppt(input_pptx_path)
