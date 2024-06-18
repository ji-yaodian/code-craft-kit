import os
from pptx import Presentation


def merge_ppt(input_pptx_paths, output_pptx_path):
    """
    将多个PPT文件合并成一个PPT文件。

    Args:
        input_pptx_paths (list[str]): 输入PPT文件路径列表。
        output_pptx_path (str): 输出PPT文件路径。
    """

    # 创建输出演示文稿
    output_presentation = Presentation()
    print(input_pptx_paths)
    print(output_pptx_path)
    # 合并幻灯片
    for input_pptx_path in input_pptx_paths:
        input_presentation = Presentation(input_pptx_path)
        for slide in input_presentation.slides:
            output_presentation.slides.add_slide(slide)

    # 保存合并后的PPT文件
    output_presentation.save(output_pptx_path)
def convert_slide_to_layout(slide, presentation):
    """
    Converts a Slide to a SlideLayout.

    Args:
        slide (pptx.slide.Slide): The Slide object to convert.
        presentation (pptx.presentation.Presentation): The presentation containing the Slide.
    """

    # Extract layout information
    slide_master = presentation.slide_master

    # Get SlideLayout from Slide
    slide_layout = slide.slide_layout

    # Create new SlideLayout
    new_slide_layout = slide_master.slide_layouts[1]

    # Set layout properties (example)
    new_slide_layout.title = slide_layout.title
    new_slide_layout.background = slide_layout.background

    # Apply layout to a new Slide
    new_slide = presentation.slides.add_slide()
    new_slide.slide_layout = new_slide_layout

    return new_slide_layout

def merge_ppt2(input_pptx_paths, output_pptx_path):
    """
    Merges multiple PPT files into a single PPT file.

    Args:
        input_pptx_paths (list[str]): List of input PPT file paths.
        output_pptx_path (str): Output PPT file path.
    """

    output_presentation = Presentation()

    for input_pptx_path in input_pptx_paths:
        input_presentation = Presentation(input_pptx_path)
        for slide in input_presentation.slides:
            # Use split_layout to get SlideLayout object
            slide_layout = convert_slide_to_layout(slide, output_presentation)

            # Clone placeholder shapes from SlideLayout to output_presentation
            for placeholder in slide_layout.iter_cloneable_placeholders():
                output_presentation.slides.add_slide(slide_layout).shapes.add_shape(
                    placeholder.shape_type,
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                )

    output_presentation.save(output_pptx_path)

if __name__ == "__main__":
    import sys

    # 输入PPT文件路径列表

    merge_ppt2(input_pptx_paths, output_pptx_path)
